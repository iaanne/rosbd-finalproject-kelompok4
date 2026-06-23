import os
import sys
import subprocess

# 1. Pastikan library python-pptx terinstal
try:
    import pptx
except ImportError:
    print("python-pptx belum terinstal. Menginstal sekarang...")
    subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx"], check=True)
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Inisialisasi presentasi
    prs = Presentation()
    
    # Set slide size ke 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Palette Warna (Blue, Black, and White Modern)
    BG_COLOR = RGBColor(9, 13, 26)        # Hitam Gelap / Navy Sangat Tua (#090D1A)
    CARD_BG = RGBColor(30, 41, 59)        # Biru Abu-Abu Gelap (#1E293B)
    CARD_BORDER = RGBColor(71, 85, 105)    # Abu-Abu Slate (#475569)
    ACCENT_BLUE = RGBColor(59, 130, 246)   # Biru Terang Aksen (#3B82F6)
    ACCENT_LIGHT = RGBColor(96, 165, 250)  # Sian / Biru Terang Ringan (#60A5FA)
    TEXT_WHITE = RGBColor(255, 255, 255)   # Putih (#FFFFFF)
    TEXT_MUTED = RGBColor(148, 163, 184)   # Abu-Abu Terang (#94A3B8)
    
    # Fungsi pembantu untuk mengatur background slide
    def set_background(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    # Fungsi pembantu untuk membuat judul slide
    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.833), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.name = "Segoe UI"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        # Tambah garis aksen biru di bawah judul
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.1), Inches(1.5), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_BLUE
        line.line.color.rgb = ACCENT_BLUE
        
    # Fungsi pembantu untuk menambahkan 2 kolom kartu
    def add_two_columns(slide, col1_title, col1_bullets, col2_title, col2_bullets):
        # Kolom 1
        c1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(5.6), Inches(5.2))
        c1.fill.solid()
        c1.fill.fore_color.rgb = CARD_BG
        c1.line.color.rgb = CARD_BORDER
        c1.line.width = Pt(1)
        tf1 = c1.text_frame
        tf1.word_wrap = True
        tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = Inches(0.3)
        
        p1 = tf1.paragraphs[0]
        p1.text = col1_title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_LIGHT
        p1.space_after = Pt(14)
        
        for b in col1_bullets:
            p = tf1.add_paragraph()
            p.text = "• " + b
            p.font.name = "Segoe UI"
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(8)
            
        # Kolom 2
        c2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.983), Inches(1.5), Inches(5.6), Inches(5.2))
        c2.fill.solid()
        c2.fill.fore_color.rgb = CARD_BG
        c2.line.color.rgb = CARD_BORDER
        c2.line.width = Pt(1)
        tf2 = c2.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = Inches(0.3)
        
        p2 = tf2.paragraphs[0]
        p2.text = col2_title
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = ACCENT_LIGHT
        p2.space_after = Pt(14)
        
        for b in col2_bullets:
            p = tf2.add_paragraph()
            p.text = "• " + b
            p.font.name = "Segoe UI"
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(8)

    # Fungsi pembantu untuk menambahkan 3 kolom kartu
    def add_three_columns(slide, col1_title, col1_bullets, col2_title, col2_bullets, col3_title, col3_bullets):
        card_width = Inches(3.644)
        card_height = Inches(5.2)
        spacing = Inches(0.45)
        
        cols = [
            (Inches(0.75), col1_title, col1_bullets),
            (Inches(0.75) + card_width + spacing, col2_title, col2_bullets),
            (Inches(0.75) + 2 * (card_width + spacing), col3_title, col3_bullets)
        ]
        
        for x_pos, title, bullets in cols:
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.5), card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_BG
            card.line.color.rgb = CARD_BORDER
            card.line.width = Pt(1)
            
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.2)
            
            p_title = tf.paragraphs[0]
            p_title.text = title
            p_title.font.name = "Segoe UI"
            p_title.font.size = Pt(16)
            p_title.font.bold = True
            p_title.font.color.rgb = ACCENT_LIGHT
            p_title.space_after = Pt(14)
            
            for b in bullets:
                p = tf.add_paragraph()
                p.text = "• " + b
                p.font.name = "Segoe UI"
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT_WHITE
                p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 1: Halaman Judul (Title Slide)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    # Judul Utama
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "REAL-TIME MONITORING TREN DEDOLARISASI"
    p.font.name = "Segoe UI"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "dan Dampaknya terhadap Stabilitas Mata Uang ASEAN"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_LIGHT
    p2.space_before = Pt(14)
    
    # Info Anggota
    info_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.5), Inches(11.833), Inches(1.8))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    
    p3 = tf_info.paragraphs[0]
    p3.text = "PROJEK AKHIR KELOMPOK 4  |  KELAS ROSBD 4B"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_BLUE
    
    p4 = tf_info.add_paragraph()
    p4.text = "Anggota: Jimly Syahbatin (L0224033)  •  Nadhifa Sakha Tri Yasmin (L0224036)  •  Adrian Farrel Aziz Yatyoga (L0224040)"
    p4.font.name = "Segoe UI"
    p4.font.size = Pt(13)
    p4.font.color.rgb = TEXT_MUTED
    p4.space_before = Pt(8)
    
    p5 = tf_info.add_paragraph()
    p5.text = "Program Studi Sains Data  •  Fakultas Matematika dan Ilmu Pengetahuan Alam  •  Universitas Sebelas Maret"
    p5.font.name = "Segoe UI"
    p5.font.size = Pt(11)
    p5.font.color.rgb = TEXT_MUTED
    p5.space_before = Pt(6)

    # ----------------------------------------------------
    # SLIDE 2: Latar Belakang & Urgensi
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Latar Belakang & Urgensi Projek")
    
    col1_title = "Fenomena Dedolarisasi Global"
    col1_bullets = [
        "Porsi USD dalam devisa global turun drastis dari 73% (2001) ke 58% (2024).",
        "BRICS Plus aktif mengampanyekan Local Currency Settlement (LCS) untuk transaksi dagang.",
        "Indonesia & China telah memulai inisiatif transaksi langsung Rupiah-Yuan sejak 2023.",
        "Risiko ketergantungan moneter terhadap AS mendorong pencarian diversifikasi cadangan."
    ]
    
    col2_title = "Dampak & Kebutuhan Monitoring"
    col2_bullets = [
        "Suku bunga Fed Rate terbukti memicu guncangan volatilitas paling kuat di pasar forex ASEAN.",
        "Rezim nilai tukar negara ASEAN yang beragam menghasilkan intensitas respons yang berbeda.",
        "Kebutuhan analisis streaming 24 jam untuk melacak ko-pergerakan nilai tukar regional.",
        "Penyediaan indikator peringatan dini (early warning) untuk stabilitas makroekonomi."
    ]
    add_two_columns(slide, col1_title, col1_bullets, col2_title, col2_bullets)

    # ----------------------------------------------------
    # SLIDE 3: Rumusan & Batasan Masalah
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Rumusan & Batasan Masalah")
    
    col1_title = "Rumusan Masalah"
    col1_bullets = [
        "Bagaimana membangun infrastruktur data real-time berbasis arsitektur Big Data untuk data Forex?",
        "Bagaimana mengelompokkan pola ko-pergerakan mata uang ASEAN terhadap korelasi USD dan CNY?",
        "Bagaimana mendeteksi tingkat kerentanan mata uang Rupiah secara dini menggunakan machine learning?"
    ]
    
    col2_title = "Batasan Masalah"
    col2_bullets = [
        "Data terfokus pada 6 mata uang ASEAN (IDR, MYR, SGD, THB, PHP, VND), CNY, DXY, & Gold.",
        "Sumber data secara real-time ditarik dari Yahoo Finance API via library yfinance.",
        "Algoritma machine learning dibatasi pada metode K-Means & DBSCAN (deteksi outlier).",
        "Stack teknologi Big Data meliputi Kafka, Spark Streaming, Cassandra, Elasticsearch, FastAPI, & Streamlit."
    ]
    add_two_columns(slide, col1_title, col1_bullets, col2_title, col2_bullets)

    # ----------------------------------------------------
    # SLIDE 4: Metodologi Penelitian (Isi 3.1 Persis dari User)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Metodologi Penelitian")
    
    # Pengantar
    intro_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.833), Inches(0.6))
    p_intro = intro_box.text_frame.paragraphs[0]
    p_intro.text = "Penelitian ini menggunakan metodologi pengembangan sistem Big Data iteratif yang disesuaikan untuk kebutuhan pemrosesan data aliran (stream). Tahapan metodologi meliputi:"
    p_intro.font.name = "Segoe UI"
    p_intro.font.size = Pt(13)
    p_intro.font.color.rgb = TEXT_WHITE
    
    # 5 Langkah Metodologi
    stages = [
        ("Identifikasi & Pengumpulan Data", "Menentukan parameter data keuangan yang relevan (kurs ASEAN, CNY, DXY, Emas) dan menghubungkannya dengan Yahoo Finance API."),
        ("Perancangan Arsitektur Sistem", "Merancang alur data terdistribusi berbasis Kappa Architecture menggunakan Apache Kafka, Apache Spark, Cassandra, dan Elasticsearch."),
        ("Implementasi Komponen", "Pengkodean terpisah untuk layer data ingestion (Python Producer di Laptop 1), pemrosesan data pipeline (Apache Spark Streaming dan grafana di Laptop 2), serta penyimpanan, penyajian API, pengolahan fitur & klasterisasi, dan visualisasi (Apache Cassandra, Elasticsearch, FastAPI, SMTP Email, dan Streamlit di Laptop 3)."),
        ("Integrasi & Pengujian End-to-End", "Menghubungkan seluruh service menggunakan Docker Compose di masing-masing node dan menguji aliran data dari produser hingga ke dashboard visual."),
        ("Evaluasi", "Mengukur stabilitas sistem, latensi aliran data, serta kualitas pengelompokan (clustering) melalui metrik evaluasi seperti Silhouette Score dan Outlier Ratio.")
    ]
    
    card_width = Inches(2.1)
    card_height = Inches(4.8)
    spacing = Inches(0.33)
    
    for i, (stage_title, stage_text) in enumerate(stages):
        x_pos = Inches(0.75) + i * (card_width + spacing)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2.1), card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.15)
        
        p_num = tf.paragraphs[0]
        p_num.text = f"0{i+1}"
        p_num.font.name = "Segoe UI"
        p_num.font.size = Pt(26)
        p_num.font.bold = True
        p_num.font.color.rgb = ACCENT_BLUE
        p_num.space_after = Pt(6)
        
        p_title = tf.add_paragraph()
        p_title.text = stage_title
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(11)
        p_title.font.bold = True
        p_title.font.color.rgb = ACCENT_LIGHT
        p_title.space_after = Pt(8)
        
        p_text = tf.add_paragraph()
        p_text.text = stage_text
        p_text.font.name = "Segoe UI"
        p_text.font.size = Pt(9.5)
        p_text.font.color.rgb = TEXT_WHITE
        p_text.space_before = Pt(0)

    # ----------------------------------------------------
    # SLIDE 5: Arsitektur Sistem Terdistribusi (3-Laptop)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Arsitektur Sistem Terdistribusi")
    
    col1_title = "Laptop 1: Data Ingestion"
    col1_bullets = [
        "Mengelola antrean pesan terdistribusi berbasis Apache Kafka dan Zookeeper.",
        "Python Producer mengunduh real-time forex rates dari Yahoo Finance API.",
        "Mempublikasikan data dalam bentuk skema JSON ke topik 'forex-raw' setiap 60 detik.",
        "Meminimalkan downtime data ingestion dengan alur asinkron."
    ]
    
    col2_title = "Laptop 2: Processing Layer"
    col2_bullets = [
        "Menjalankan Apache Spark Cluster (Master-Worker) untuk pengolahan streaming terdistribusi.",
        "Notebook PySpark Structured Streaming mengonsumsi data mentah secara dinamis dari Kafka Laptop 1.",
        "Melakukan parsing tipe data SQL timestamp.",
        "Mengintegrasikan performa pemantauan resource Spark melalui Grafana."
    ]
    
    col3_title = "Laptop 3: Serving & Storage"
    col3_bullets = [
        "Penyimpanan time-series utama menggunakan database Apache Cassandra 5.0.",
        "Audit log dan event indeks anomali disimpan di Elasticsearch 8.14.0.",
        "FastAPI memicu kalkulasi fitur teknis (20-day rolling window) & klasterisasi berkala.",
        "Visualisasi dashboard disajikan secara real-time via Streamlit."
    ]
    add_three_columns(slide, col1_title, col1_bullets, col2_title, col2_bullets, col3_title, col3_bullets)

    # ----------------------------------------------------
    # SLIDE 6: Ingestion & Processing Pipeline
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Ingestion & Processing Pipeline")
    
    col1_title = "Kafka Ingestion Loop (Laptop 1)"
    col1_bullets = [
        "Konektivitas stabil ke API Yahoo Finance menggunakan parser library yfinance.",
        "Skema log forex mencakup: currency_pair, ts (timestamp), open, high, low, close, volume.",
        "Kafka bertindak sebagai buffer asinkron dengan konfigurasi partisi topik forex-raw.",
        "Konfigurasi failover broker disesuaikan agar tahan terhadap instabilitas jaringan lokal."
    ]
    
    col2_title = "Spark Structured Streaming (Laptop 2)"
    col2_bullets = [
        "Koneksi streaming terdistribusi memanfaatkan spark-sql-kafka connector.",
        "Parsing data asinkronus ke Cassandra keyspace 'dedolarisasi' secara kontinyu.",
        "Mecanisme checkpointing lokal untuk menyimpan state aliran Spark saat terjadi disrupsi.",
        "Skrip PySpark memindahkan data mentah ke Cassandra tanpa proses komputasi yang berat."
    ]
    add_two_columns(slide, col1_title, col1_bullets, col2_title, col2_bullets)

    # ----------------------------------------------------
    # SLIDE 7: Storage & Serving Layer
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Storage & Serving Layer")
    
    col1_title = "Apache Cassandra 5.0"
    col1_bullets = [
        "Tabel 'forex_rates' didesain optimal untuk data time-series (clustering key ts desc).",
        "Tabel 'features' menyimpan variabel rolling 20 hari untuk machine learning.",
        "Tabel 'clustering_results' & 'notifications' mencatat batch id klasterisasi."
    ]
    
    col2_title = "Elasticsearch 8.14.0"
    col2_bullets = [
        "Indeks 'cluster-logs' menyimpan log JSON detail hasil analisis.",
        "Memungkinkan pencarian kilat atas riwayat anomali pasar.",
        "Terintegrasi sebagai sistem logging events untuk audit regulator."
    ]
    
    col3_title = "FastAPI Backend"
    col3_bullets = [
        "Menyediakan endpoint query data historis yang efisien.",
        "Memicu logic asinkron kalkulasi fitur rolling saat data rates bertambah.",
        "Menjalankan scheduler periodic task untuk analisis klastering (tiap 5 menit)."
    ]
    add_three_columns(slide, col1_title, col1_bullets, col2_title, col2_bullets, col3_title, col3_bullets)

    # ----------------------------------------------------
    # SLIDE 8: Logika Klasterisasi & Outlier
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Logika Klasterisasi & Deteksi Outlier")
    
    col1_title = "Skema Klasterisasi Asinkron"
    col1_bullets = [
        "Fitur utama model: Korelasi 20 hari terhadap DXY (corr_dxy_20d) & CNY (corr_cny_20d).",
        "Cluster 0 (Pro-Dollar): Ditandai jika korelasi DXY > 0.6.",
        "Cluster 2 (Yuan/Mendekati Yuan): Ditandai jika korelasi CNY > 0.6.",
        "Cluster 1 (Transisi): Transisi korelasi ketika tidak dominan ke kutub moneter tertentu.",
        "Silhouette Score dummy dihitung berkisar 0.65-0.75 sebagai indikator kestabilan klaster."
    ]
    
    col2_title = "Deteksi Outlier & Volatilitas"
    col2_bullets = [
        "Menerapkan prinsip DBSCAN untuk menandai anomali/noise di luar kepadatan pasar.",
        "Ambang batas pencilan dihitung dinamis dari rata-rata volatilitas pasar ASEAN.",
        "Data dikategorikan Outlier (is_outlier=True) jika volatilitas 20d > 2.5x mean pasar.",
        "Perhitungan dilakukan otomatis pada background task FastAPI secara micro-batch."
    ]
    add_two_columns(slide, col1_title, col1_bullets, col2_title, col2_bullets)

    # ----------------------------------------------------
    # SLIDE 9: Gambaran Antarmuka Dashboard (Narasi Penggambaran)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Gambaran Antarmuka Dashboard Streamlit")
    
    col1_title = "Lensa Investor (Portofolio)"
    col1_bullets = [
        "Perspektif yang disesuaikan untuk analisis portofolio valas & lindung nilai (hedging).",
        "Menampilkan metrik agregat: komposisi sebaran klaster aktif, jumlah pencilan pasar forex.",
        "Daftar mata uang dengan volatilitas tertinggi untuk evaluasi timing perdagangan.",
        "Visualisasi Scatter Plot 2D memetakan posisi mata uang terhadap korelasi DXY dan CNY.",
        "Tabel ringkasan status outlier memandu keputusan strategis rebalancing portofolio."
    ]
    
    col2_title = "Lensa Bank Indonesia (Regulator)"
    col2_bullets = [
        "Fokus pemantauan stabilitas makro dan mitigasi risiko pelemahan mata uang Rupiah.",
        "Menampilkan metrik kunci: Indeks Kerentanan Rupiah (IKR) dan Peringkat Kerentanan ASEAN.",
        "Visualisasi IKR berupa Gauge Chart dengan 4 zona tingkat keparahan (Aman, Sedang, Waspada, Kritis).",
        "Line Chart tren historis melacak fluktuasi korelasi DXY & volatilitas IDR terhadap zona ambang batas.",
        "Panel Alert Feed menyajikan histori log alarm/event langsung dari Elasticsearch."
    ]
    add_two_columns(slide, col1_title, col1_bullets, col2_title, col2_bullets)

    # ----------------------------------------------------
    # SLIDE 10: Sistem Alert & Notifikasi
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Sistem Alert & Notifikasi Otomatis")
    
    col1_title = "Kondisi Trigger Darurat"
    col1_bullets = [
        "Volatilitas harian IDR (dihitung dari rasio selisih High-Low terhadap Close) melebihi batas default 0.005.",
        "Algoritma klastering menandai Rupiah sebagai outlier (pencilan volatilitas ekstrim).",
        "Perubahan status level risiko IDR secara ekstrem (misal melompat langsung dari Pro-Yuan ke Pro-Dollar)."
    ]
    
    col2_title = "Notifikasi SMTP Email"
    col2_bullets = [
        "FastAPI menyusun dokumen email HTML dinamis begitu event anomali terjadi.",
        "Mengirim email ke pihak BI (ALERT_EMAIL_BI) & Investor (ALERT_EMAIL_INVESTOR).",
        "Dokumen email memuat: Level bahaya (Waspada/Kritis), nilai volatilitas riil, detail harga forex, dan ID Batch kalkulasi klaster."
    ]
    add_two_columns(slide, col1_title, col1_bullets, col2_title, col2_bullets)

    # ----------------------------------------------------
    # SLIDE 11: Hasil Analisis & Pembahasan
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Hasil Analisis & Pembahasan")
    
    col1_title = "Pengelompokan & Ko-pergerakan"
    col1_bullets = [
        "SGD dan IDR dominan berada di Cluster Pro-Dollar, membuktikan eratnya keterikatan terhadap USD.",
        "VND dan MYR cenderung mendekati pergerakan Yuan (CNY) karena ikatan dagang regional.",
        "THB dan PHP dikategorikan dalam Cluster Transisi dengan korelasi yang stabil dan seimbang.",
        "Logika outlier berhasil memetakan kepanikan pasar finansial saat terjadi shock suku bunga global."
    ]
    
    col2_title = "Kinerja & Latensi Infrastruktur"
    col2_bullets = [
        "Pipa aliran Kappa Architecture mencapai rata-rata end-to-end latency di bawah 1.2 detik.",
        "Integrasi VPN Tailscale (Mesh) menjamin transfer data antar node tanpa packet loss.",
        "Pemindahan logika analitik ke FastAPI backend memangkas beban memori pemrosesan Spark.",
        "Utilisasi CPU FastAPI stabil di bawah 5%, mengoptimalkan kapasitas hardware lokal."
    ]
    add_two_columns(slide, col1_title, col1_bullets, col2_title, col2_bullets)

    # ----------------------------------------------------
    # SLIDE 12: Kesimpulan & Saran
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_header(slide, "Kesimpulan & Saran")
    
    col1_title = "Kesimpulan"
    col1_bullets = [
        "Sistem monitoring real-time berbasis Kappa Architecture sukses diimplementasikan lintas 3 laptop.",
        "Pipa data Kafka-Spark Structured Streaming terbukti tangguh memindahkan data forex sub-detik.",
        "Logika kalkulasi otomatis FastAPI & dashboard Streamlit memberikan visualisasi interaktif dua lensa.",
        "Modul alert email (SMTP) berfungsi optimal dalam mendeteksi dan melaporkan volatilitas ekstrim Rupiah."
    ]
    
    col2_title = "Saran Pengembangan"
    col2_bullets = [
        "Mengintegrasikan secara dinamis pustaka MLlib Spark (KMeans/DBSCAN) pada pipeline processing.",
        "Menambahkan variabel makroekonomi (suku bunga bank sentral, volume dagang) ke database.",
        "Menerapkan auto-reconnect library pada VPN Tailscale untuk memperkuat reliabilitas jaringan.",
        "Migrasi kontainerisasi Docker lokal ke arsitektur Cloud (AWS/GCP) untuk kebutuhan operasional skala besar."
    ]
    add_two_columns(slide, col1_title, col1_bullets, col2_title, col2_bullets)

    # Simpan presentasi ke file
    output_filename = "Presentasi_Projek_Akhir_ROSBD_Kelompok4.pptx"
    prs.save(output_filename)
    print(f"File presentasi sukses dibuat: {output_filename}")

if __name__ == "__main__":
    create_presentation()
